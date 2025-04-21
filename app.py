from flask import Flask, request, jsonify, send_file
from werkzeug.utils import secure_filename
from pptx import Presentation
import os
from evaluator import evaluate_pitch_deck
from report_generator import generate_pdf_report

UPLOAD_FOLDER = 'uploads'
REPORT_FOLDER = 'reports'

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50 MB

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(REPORT_FOLDER, exist_ok=True)

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files['file']
    if file.filename == '' or not file.filename.endswith(('.pptx', '.ppt')):
        return jsonify({"error": "Invalid file format"}), 400

    filename = secure_filename(file.filename)
    path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(path)

    try:
        pres = Presentation(path)
        content = "\n".join([
            shape.text for slide in pres.slides for shape in slide.shapes if hasattr(shape, "text")
        ])
    except Exception as e:
        return jsonify({"error": f"Failed to extract text: {str(e)}"}), 500

    evaluation_result = evaluate_pitch_deck(content)

    timestamp_clean = evaluation_result['timestamp'].replace(':', '').replace(' ', '_')
    pdf_filename = f"Investment_Thesis_{evaluation_result['startup_name']}_{timestamp_clean}.pdf"
    pdf_path = os.path.join(REPORT_FOLDER, pdf_filename)

    generate_pdf_report(evaluation_result, pdf_path)

    return send_file(pdf_path, as_attachment=True)

@app.route('/')
def home():
    return "Automated Investment Thesis Generator API is running."


if __name__ == '__main__':
    app.run(debug=True)
